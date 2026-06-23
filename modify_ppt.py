import collections.abc
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# PPTX 文件路径
ppt_path = r"c:\Users\dou12\Desktop\youziyi\youziyi-oss\assignments\游子衣答辩PPT_harmony.pptx"

def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Add a textbox to a slide with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    if color:
        p.font.color.rgb = color
    p.alignment = alignment
    return txBox

try:
    # Load PPT
    prs = Presentation(ppt_path)
    print(f"成功加载 PPT，共 {len(prs.slides)} 页幻灯片")
    
    # ============ Step 1: 修改 P1 封面 ============
    if len(prs.slides) > 0:
        slide = prs.slides[0]
        modified = False
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                full_text = para.text
                if "亲情的记录" in full_text or "数字陪伴平台" in full_text:
                    for run in para.runs:
                        run.text = ""
                    para.runs[0].text = "亲情的记录和见证之所 —— 面向3.23亿老年人的鸿蒙原生数字陪伴平台"
                    modified = True
                    break
        if modified:
            print("  P1 封面文本已修改")
        else:
            print("  P1 未找到匹配文本（可能已修改或结构不同）")

    # ============ Step 2: 修改 P3 架构图 ============
    if len(prs.slides) > 2:
        slide = prs.slides[2]
        modified = False
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                if "双端协同" in para.text:
                    for run in para.runs:
                        if "双端协同" in run.text:
                            run.text = run.text.replace("双端协同", "三端协同")
                            modified = True
                if "Web 网页端 + 微信小程序端" in para.text and "鸿蒙" not in para.text:
                    for run in para.runs:
                        if "Web 网页端 + 微信小程序端" in run.text:
                            run.text = run.text.replace(
                                "Web 网页端 + 微信小程序端",
                                "Web 网页端 + 微信小程序端 + HarmonyOS 鸿蒙端"
                            )
                            modified = True
        if modified:
            print("  P3 架构图文本已修改")
        else:
            print("  P3 未找到匹配文本")

    # ============ Step 3: 修改 P8 技术架构 ============
    if len(prs.slides) > 7:
        slide = prs.slides[7]
        modified = False
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                if "ArkTS 鸿蒙版也在同步开发" in para.text:
                    for run in para.runs:
                        if "ArkTS 鸿蒙版也在同步开发" in run.text:
                            run.text = run.text.replace(
                                'ArkTS \u9e3f\u8499\u7248\u4e5f\u5728\u540c\u6b65\u5f00\u53d1\uff0c\u8fdb\u4e00\u6b65\u62d3\u5c55\u9002\u8001\u573a\u666f\u7684\u8bbe\u5907\u8986\u76d6\u2014\u2014\u54cd\u5e94\u56fd\u5bb6\u201c\u6df1\u5316\u667a\u6167\u5065\u5eb7\u517b\u8001\u5e94\u7528\u793a\u8303\u201d\u7684\u653f\u7b56\u5bfc\u5411\u3002',
                                'ArkTS \u9e3f\u8499\u7aef\uff08\u6838\u5fc3\u843d\u5730\u70b9\uff09\u2014 \u57fa\u4e8e ArkTS + ArkUI \u58f0\u660e\u5f0f\u5f00\u53d1\uff0c\u5df2\u5b9e\u73b0 AVRecorder/AVPlayer \u8bed\u97f3\u529f\u80fd\u3001@ohos.net.http \u7f51\u7edc\u8bf7\u6c42\u3001i18n \u56fd\u9645\u8bed\u8a00\u9002\u914d\u7b49\u7cfb\u7edf\u80fd\u529b'
                            )
                            modified = True
        if modified:
            print("  P8 技术架构文本已修改")
        else:
            print("  P8 未找到匹配文本")

    # ============ Step 4: 插入 P8.5 鸿蒙原生能力深度展示 ============
    # Use a blank layout
    blank_layout = prs.slide_layouts[6]  # usually blank
    slide_8_5 = prs.slides.add_slide(blank_layout)
    
    # Title
    add_textbox(slide_8_5, Inches(0.5), Inches(0.3), Inches(9), Inches(0.8),
                "鸿蒙原生能力深度展示", font_size=32, bold=True, color=RGBColor(0x1A, 0x73, 0xE8))
    
    # Subtitle
    add_textbox(slide_8_5, Inches(0.5), Inches(1.2), Inches(9), Inches(0.5),
                "ArkTS + ArkUI 声明式开发  |  HarmonyOS NEXT 系统级原生能力  |  纯鸿蒙实现，无第三方依赖",
                font_size=14, color=RGBColor(0x66, 0x66, 0x66))
    
    # Content blocks
    content_items = [
        ("组件架构", "EntryAbility → Splash → 路由守卫 → Login / Onboarding / FamilyJoin / Main\nMain 使用 ArkUI Tabs 组件实现四 Tab 底部导航：首页（打卡）、互动（语音）、预警（健康）、我的"),
        ("语音接龙 (AVRecorder/AVPlayer)", "录制：AVRecorder + AAC 格式 + MIC 源 + 48000bps + 44100Hz\n播放：AVPlayer 在线流播放 + stateChange 事件监听\n存储：@ohos.file.fs 写入缓存 → API 上传"),
        ("路由守卫与状态管理", "AppSession 全局状态 (userId, userRole, profile)\nresolveAuthenticatedRoute 检查 Onboarding + 家庭连接\n三端路由逻辑完全一致"),
        ("网络与数据层", "@ohos.net.http 原生请求 → AppRepository Mock/远程双轨切换\n@ohos.promptAction 原生 Toast 提示"),
    ]
    
    y_pos = 2.0
    for title_text, desc_text in content_items:
        add_textbox(slide_8_5, Inches(0.7), Inches(y_pos), Inches(8.5), Inches(0.35),
                    title_text, font_size=16, bold=True, color=RGBColor(0x1A, 0x73, 0xE8))
        y_pos += 0.35
        add_textbox(slide_8_5, Inches(0.9), Inches(y_pos), Inches(8.3), Inches(0.7),
                    desc_text, font_size=13, color=RGBColor(0x33, 0x33, 0x33))
        y_pos += 0.75
    
    print("  P8.5 鸿蒙原生能力深度展示 - 已插入")

    # ============ Step 5: 插入 P8.6 国际语言适配 i18n ============
    slide_8_6 = prs.slides.add_slide(blank_layout)
    
    # Title
    add_textbox(slide_8_6, Inches(0.5), Inches(0.3), Inches(9), Inches(0.8),
                "国际语言适配（i18n）—— 技术亮点专章", font_size=32, bold=True, color=RGBColor(0x1A, 0x73, 0xE8))
    
    # Subtitle
    add_textbox(slide_8_6, Inches(0.5), Inches(1.2), Inches(9), Inches(0.5),
                "HarmonyOS 原生 i18n 资源管理  |  零运行时开销  |  适老化 + 国际化双重价值",
                font_size=14, color=RGBColor(0x66, 0x66, 0x66))
    
    # Content blocks
    i18n_items = [
        ("三层资源目录结构", 
         "AppScope/resources/ 与 entry/src/main/resources/ 下均建立：\n"
         "  base/element/string.json  → 默认语言资源（中文）\n"
         "  zh_CN/element/string.json → 简体中文专有资源\n"
         "  en_US/element/string.json → 英文资源"),
        ("$string 引用机制", 
         "app.json5:  \"label\": \"$string:app_name\"  → 中文「游子衣」/ 英文「YouZiYi」\n"
         "module.json5: \"description\": \"$string:module_desc\"\n"
         "系统根据设备语言自动选择，应用层无需 if-else 判断"),
        ("多语言实例对照", 
         "module_desc:\n"
         "  中文 → 游子衣鸿蒙端前端骨架\n"
         "  英文 → YouZiYi HarmonyOS Client Skeleton\n"
         "microphone_reason:\n"
         "  中文 → 用于录制家庭语音消息，实现亲情语音接龙功能\n"
         "  英文 → Used to record family voice messages for the voice relay feature"),
        ("核心价值", 
         "1. 零运行时开销：资源选择由系统层完成\n"
         "2. 可扩展性强：新增语言只需添加资源目录（如 ja_JP/），无需改业务代码\n"
         "3. 适老化+国际化双重价值：海外华人家庭——父母中文 / 子女英文——无缝适配"),
    ]
    
    y_pos = 2.0
    for title_text, desc_text in i18n_items:
        add_textbox(slide_8_6, Inches(0.7), Inches(y_pos), Inches(8.5), Inches(0.35),
                    title_text, font_size=16, bold=True, color=RGBColor(0x1A, 0x73, 0xE8))
        y_pos += 0.35
        lines = desc_text.count('\n') + 1
        box_height = max(0.5, lines * 0.22)
        add_textbox(slide_8_6, Inches(0.9), Inches(y_pos), Inches(8.3), Inches(box_height),
                    desc_text, font_size=13, color=RGBColor(0x33, 0x33, 0x33))
        y_pos += box_height + 0.15
    
    print("  P8.6 国际语言适配 i18n - 已插入")

    # ============ Step 6: Save ============
    output_path = ppt_path.replace(".pptx", "_modified.pptx")
    prs.save(output_path)
    print(f"\n修改完成！已另存为: {output_path}")
    print(f"修改后共 {len(prs.slides)} 页幻灯片")

except Exception as e:
    import traceback
    print(f"修改失败: {e}")
    traceback.print_exc()
